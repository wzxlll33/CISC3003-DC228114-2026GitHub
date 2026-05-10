from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_DIR = Path(__file__).resolve().parents[1].parent
PPT_PATH = PROJECT_DIR / "Individual Digital Story.pptx"
BACKUP_PATH = PROJECT_DIR / "Individual Digital Story.team-template-backup.pptx"


SLIDES = [
    {
        "eyebrow": "CISC3003 Individual Digital Story",
        "title": "Taste of Macau",
        "subtitle": "Catalogue, Search and Favorites Contribution",
        "bullets": [
            "Leong Chi Long - DC227153 - Team 04",
            "Silent recording prepared for later voice narration",
            "Demo scope: catalogue rendering, category filtering, i18n, guest favorites, login sync, search history",
        ],
    },
    {
        "eyebrow": "Personal Introduction",
        "title": "My Role",
        "subtitle": "Front-End Developer: discovery and saved-food experience",
        "bullets": [
            "Built the restaurant-first catalogue used on the Explore page",
            "Implemented category filtering, debounced search, API-backed fallback and search-history trigger",
            "Implemented guest favorites with LocalStorage and merge-to-account behavior after login",
            "Worked with the dashboard pages to prove favorites and search history are persisted",
        ],
    },
    {
        "eyebrow": "Demo Roadmap",
        "title": "One Linear Flow",
        "subtitle": "No repeated back-and-forth: each step proves one contribution",
        "bullets": [
            "Start as a guest and show the catalogue baseline",
            "Open tools, filter by Desserts, and verify the result count changes",
            "Switch Chinese, English and Portuguese to prove i18n behavior",
            "Save Portuguese Egg Tart as a guest, then log in and show the synced favorite",
            "Search for Lord, press Enter, then confirm the query appears in Search History",
            "Finish with code walkthrough, challenges, solutions and learning",
        ],
    },
    {
        "eyebrow": "System Context",
        "title": "Restaurant-First Catalogue",
        "subtitle": "The Explore page renders restaurants while using food categories as filters",
        "bullets": [
            "Dataset shown in the live app: 16 restaurants and 49 foods",
            "Restaurant cards include area, category, rating, signature dish cues and detail links",
            "Category filters are generated from the food category data and applied to restaurant cards",
            "The map and result sheet stay connected so discovery is visual and browsable",
        ],
    },
    {
        "eyebrow": "Live Demo 01",
        "title": "Catalogue Rendering",
        "subtitle": "Baseline browsing before any filter",
        "bullets": [
            "Open /explore as a guest",
            "Show the initial result count: 16 restaurant cards",
            "Point out rating, area, category badge, signature dishes and detail actions",
            "Use this as the baseline for the filtering test",
        ],
    },
    {
        "eyebrow": "Live Demo 02",
        "title": "Category Filtering",
        "subtitle": "A visible click path proves the filter is applied",
        "bullets": [
            "Open the tools button, then the filter panel",
            "Click Desserts and apply the filter",
            "Verify the result count changes from 16 to 7 restaurants",
            "Explain that this prevents the previous issue where filtering was mentioned but not clearly tested",
        ],
    },
    {
        "eyebrow": "Live Demo 03",
        "title": "Trilingual i18n",
        "subtitle": "Language switching is visible through UI labels and placeholders",
        "bullets": [
            "Open the language menu",
            "Switch from Chinese to English, then Portuguese, then back to English",
            "Check navigation labels and search placeholder after each reload",
            "Keep English for the later login, dashboard and code explanation",
        ],
    },
    {
        "eyebrow": "Live Demo 04",
        "title": "Guest Favorites",
        "subtitle": "Save a dish before login",
        "bullets": [
            "Open Portuguese Egg Tart detail page",
            "Click Add to Favorites as a guest",
            "The button changes to Remove from Favorites and LocalStorage stores food id 8",
            "Open guest favorites view to show matching restaurants for that saved dish",
        ],
    },
    {
        "eyebrow": "Live Demo 05",
        "title": "Login Sync",
        "subtitle": "Guest favorite becomes account data",
        "bullets": [
            "Log in with the verified Leong demo account",
            "FavoriteStore syncs guest LocalStorage IDs to the server",
            "Open Dashboard > Favorites",
            "Show Portuguese Egg Tart as a persisted favorite item",
        ],
    },
    {
        "eyebrow": "Live Demo 06",
        "title": "Search and History",
        "subtitle": "Search is only logged after Enter confirms the query",
        "bullets": [
            "Return to Explore as a logged-in user",
            "Search for Lord, a keyword that returns Lord Stow's Bakery",
            "Press Enter so commitHistory becomes true",
            "Open Dashboard > Search History and show Lord with one result",
        ],
    },
    {
        "eyebrow": "Code Walkthrough 01",
        "title": "Catalogue and Filter Logic",
        "subtitle": "public/assets/js/restaurant-catalog.js",
        "bullets": [
            "bindFilterPanel stores a draft category when a filter button is clicked",
            "applyDraftFilters updates active filters and rerenders cards",
            "applyFilters combines category, service mode and search term",
            "renderCards and syncAllUi keep count, result sheet and map aligned",
        ],
    },
    {
        "eyebrow": "Code Walkthrough 02",
        "title": "Debounced Search",
        "subtitle": "Input gives fast feedback; Enter commits history",
        "bullets": [
            "input event waits 220 ms before running search",
            "Enter clears the debounce timer and calls search with commitHistory true",
            "Local matches are tried first; API-backed restaurant search handles fallback",
            "logSearchHistory posts query, filters and result count only for logged-in users",
        ],
    },
    {
        "eyebrow": "Code Walkthrough 03",
        "title": "Favorites Sync",
        "subtitle": "Guest LocalStorage plus authenticated API persistence",
        "bullets": [
            "Guest IDs are stored in taste-of-macau:guest-favorites",
            "Logged-in toggles use /api/favorites/{foodId}",
            "On DOMContentLoaded, FavoriteStore.syncLocalToServer posts guest IDs",
            "The backend validates food IDs and writes them through FavoriteRepository",
        ],
    },
    {
        "eyebrow": "Challenge and Solution",
        "title": "What Was Hard",
        "subtitle": "The work was not just buttons; it was state consistency",
        "bullets": [
            "Challenge: guest and logged-in favorites used different storage layers",
            "Solution: normalize food IDs and make syncLocalToServer the single merge point",
            "Challenge: search felt responsive but history should not log every keystroke",
            "Solution: debounce typing, but only commit history on Enter or confirmed search",
            "Challenge: catalogue filters had to update cards, counts and map together",
        ],
    },
    {
        "eyebrow": "Learning",
        "title": "What I Learned",
        "subtitle": "Front-end state needs a contract with backend persistence",
        "bullets": [
            "Designing a clear user journey makes testing easier than checking isolated buttons",
            "Local UI state, LocalStorage and server state must be synchronized deliberately",
            "Small details such as Enter-to-log-history affect whether a demo proves the real logic",
            "Using stable data attributes made the UI testable and easier to record",
        ],
    },
    {
        "eyebrow": "Closing",
        "title": "Contribution Summary",
        "subtitle": "Leong Chi Long - DC227153",
        "bullets": [
            "Catalogue rendering and category filtering",
            "Debounced search with API-backed results and history logging",
            "Guest favorites, account sync and dashboard evidence",
            "Challenge, solution and learning are shown after the live demo and code walkthrough",
        ],
    },
]


def clear_slide_text(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            shape.text_frame.clear()


def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False, color=(35, 43, 46), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_bullets(slide, bullets, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(18 if len(bullets) <= 4 else 16)
        paragraph.font.color.rgb = RGBColor(35, 43, 46)
        paragraph.space_after = Pt(8)
        paragraph.line_spacing = 1.08
    return box


def add_accent(slide, index):
    palette = [
        RGBColor(15, 78, 150),
        RGBColor(190, 100, 72),
        RGBColor(212, 168, 67),
        RGBColor(42, 116, 91),
    ]
    color = palette[index % len(palette)]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.85), Inches(0.35), Inches(1.2), Inches(0.34))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    frame = pill.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = f"{index + 1:02d}"
    run.font.name = "Aptos"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_footer(slide):
    add_textbox(
        slide,
        Inches(0.72),
        Inches(7.05),
        Inches(6.8),
        Inches(0.22),
        "Taste of Macau - Individual Digital Story - Team 04",
        font_size=9,
        color=(92, 101, 105),
    )


def apply_slide(slide, spec, index):
    clear_slide_text(slide)
    add_accent(slide, index)
    add_textbox(
        slide,
        Inches(0.72),
        Inches(0.55),
        Inches(7.7),
        Inches(0.3),
        spec["eyebrow"].upper(),
        font_size=11,
        bold=True,
        color=(15, 78, 150),
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.0),
        Inches(10.6),
        Inches(0.95),
        spec["title"],
        font_size=38 if index else 46,
        bold=True,
        color=(24, 32, 34),
    )
    add_textbox(
        slide,
        Inches(0.74),
        Inches(1.92),
        Inches(10.0),
        Inches(0.52),
        spec["subtitle"],
        font_size=19,
        color=(92, 75, 38),
    )
    add_bullets(slide, spec["bullets"], Inches(0.9), Inches(2.82), Inches(10.8), Inches(3.6))
    add_footer(slide)


def main():
    if not PPT_PATH.exists():
        raise FileNotFoundError(PPT_PATH)
    if not BACKUP_PATH.exists():
        copy2(PPT_PATH, BACKUP_PATH)

    presentation = Presentation(PPT_PATH)
    if len(presentation.slides) < len(SLIDES):
        raise RuntimeError(f"Template has {len(presentation.slides)} slides, but {len(SLIDES)} are required.")

    for index, spec in enumerate(SLIDES):
        apply_slide(presentation.slides[index], spec, index)

    for index in range(len(SLIDES), len(presentation.slides)):
        clear_slide_text(presentation.slides[index])

    presentation.save(PPT_PATH)
    print(f"Updated: {PPT_PATH}")
    print(f"Backup:  {BACKUP_PATH}")
    print(f"Slides:  {len(SLIDES)}")


if __name__ == "__main__":
    main()
